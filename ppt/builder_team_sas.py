from pathlib import Path
from datetime import datetime
from pptx import Presentation
from .slide_clone import clone_slide

TEMPLATE_PATH = Path("assets/templates/academic_template.pptx")

TEMPLATE_SLIDES = {
    "title": 0,
    "context": 1,
    "research_question": 2,
    "datasets": 3,
    "features": 4,
    "finding": 5,
    "validation": 8,
    "conclusion": 9,
}

FINDING_SLIDES = [5, 6, 7]


def replace_text_everywhere(slide, replacements: dict):
    for shape in slide.shapes:
        if not getattr(shape, "has_text_frame", False):
            continue
        tf = shape.text_frame
        if not tf:
            continue
        for p in tf.paragraphs:
            for run in p.runs:
                text = run.text
                if not text:
                    continue
                for k, v in replacements.items():
                    if k in text:
                        run.text = text.replace(k, v)
                        text = run.text


def build_presentation(topic: str, slides: list[dict], output_dir="output") -> str:
    template = Presentation(str(TEMPLATE_PATH))
    output = Presentation()

    finding_index = 0

    for slide_data in slides:
        slide_type = slide_data["type"]

        if slide_type == "finding":
            template_slide = template.slides[FINDING_SLIDES[finding_index % len(FINDING_SLIDES)]]
            finding_index += 1
        else:
            template_slide = template.slides[TEMPLATE_SLIDES[slide_type]]

        new_slide = clone_slide(output, template_slide)

        replace_text_everywhere(new_slide, {
            "{{TOPIC}}": topic,
            "{{TITLE}}": slide_data.get("title", ""),
            "{{SUBTITLE}}": slide_data.get("subtitle", ""),
            "{{BODY}}": "\n".join(slide_data.get("bullets", [])),
        })

    Path(output_dir).mkdir(exist_ok=True)
    filename = f"presentation_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pptx"
    output_path = Path(output_dir) / filename
    output.save(output_path)

    return str(output_path)
